"""文件上传路由 - 支持图片/PDF/Word/PPT文本提取"""
import os
from datetime import datetime
from fastapi import APIRouter, Depends, HTTPException, UploadFile, File
from auth import get_current_user
from models import User

router = APIRouter(prefix="/api/upload", tags=["文件上传"])

UPLOAD_DIR = os.path.join(os.path.dirname(os.path.dirname(__file__)), "uploads")
os.makedirs(UPLOAD_DIR, exist_ok=True)

# 支持的文件类型 (MIME + 扩展名双重检测)
ALLOWED_TYPES = {
    "image/png", "image/jpeg", "image/jpg", "image/gif", "image/webp",
    "application/pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",  # .docx
    "application/msword",  # .doc
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",  # .pptx
    "application/vnd.ms-powerpoint",  # .ppt
    "text/plain",
    "application/octet-stream",  # 有些浏览器会把 ppt/pptx 识别为这个
}

ALLOWED_EXTENSIONS = {
    ".png", ".jpg", ".jpeg", ".gif", ".webp",
    ".pdf", ".doc", ".docx", ".ppt", ".pptx", ".txt",
}

MAX_FILE_SIZE = 10 * 1024 * 1024  # 10MB


def _detect_file_type(filename: str, content_type: str) -> str:
    """通过文件扩展名补充检测真实类型"""
    ext = os.path.splitext(filename)[1].lower()
    # 如果 content_type 为通用类型，根据扩展名推断
    if content_type in ("application/octet-stream", "", None):
        ext_map = {
            ".pdf": "application/pdf",
            ".doc": "application/msword",
            ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            ".ppt": "application/vnd.ms-powerpoint",
            ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            ".png": "image/png",
            ".jpg": "image/jpeg",
            ".jpeg": "image/jpeg",
            ".gif": "image/gif",
            ".txt": "text/plain",
        }
        return ext_map.get(ext, content_type or "")
    return content_type


def extract_text_from_file(filepath: str, content_type: str, filename: str = "") -> str:
    """从文件中提取文本内容"""
    # 用扩展名补充类型检测
    real_type = _detect_file_type(filename or filepath, content_type)
    ext = os.path.splitext(filename or filepath)[1].lower()
    text = ""

    try:
        if real_type == "text/plain" or ext == ".txt":
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()

        elif real_type == "application/pdf" or ext == ".pdf":
            try:
                import PyPDF2
                with open(filepath, "rb") as f:
                    reader = PyPDF2.PdfReader(f)
                    for page in reader.pages[:20]:
                        text += page.extract_text() or ""
            except ImportError:
                text = "[PDF文件] 需要安装 PyPDF2 才能提取内容"

        elif ext in (".doc", ".docx") or "wordprocessing" in real_type or "msword" in real_type:
            try:
                import docx
                doc = docx.Document(filepath)
                for para in doc.paragraphs:
                    text += para.text + "\n"
                # 也提取表格内容
                for table in doc.tables:
                    for row in table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                        if row_text:
                            text += row_text + "\n"
            except ImportError:
                text = "[Word文档] 需要安装 python-docx 才能提取内容"
            except Exception as e:
                # .doc 旧格式 python-docx 不支持
                if ext == ".doc":
                    text = f"[Word文档] 旧版 .doc 格式，建议转换为 .docx 后重新上传"
                else:
                    text = f"[Word解析失败: {str(e)[:100]}]"

        elif ext in (".ppt", ".pptx") or "presentation" in real_type or "powerpoint" in real_type:
            try:
                from pptx import Presentation
                prs = Presentation(filepath)
                for i, slide in enumerate(prs.slides):
                    slide_text = []
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                t = para.text.strip()
                                if t:
                                    slide_text.append(t)
                        # 也提取表格
                        if shape.has_table:
                            for row in shape.table.rows:
                                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                                if row_text:
                                    slide_text.append(row_text)
                    if slide_text:
                        text += f"[第{i+1}页] " + " / ".join(slide_text) + "\n"
            except ImportError:
                text = "[PPT文件] 需要安装 python-pptx 才能提取内容"
            except Exception as e:
                # .ppt 旧格式 python-pptx 不支持
                if ext == ".ppt":
                    text = f"[PPT文件] 旧版 .ppt 格式，建议转换为 .pptx 后重新上传。提示：用WPS/PowerPoint打开后另存为.pptx格式"
                else:
                    text = f"[PPT解析失败: {str(e)[:100]}]"

        elif real_type.startswith("image/"):
            text = f"[图片文件: {os.path.basename(filepath)}]"

    except Exception as e:
        text = f"[文件解析失败: {str(e)[:100]}]"

    return text.strip()


@router.post("/file")
async def upload_file(
    file: UploadFile = File(...),
    current_user: User = Depends(get_current_user),
):
    """上传文件并提取文本内容"""
    ext = os.path.splitext(file.filename or "")[1].lower()
    real_type = _detect_file_type(file.filename or "", file.content_type)

    # 双重检测：MIME类型 或 文件扩展名
    if real_type not in ALLOWED_TYPES and ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(
            status_code=400,
            detail=f"不支持的文件类型: {file.content_type} ({ext})。支持: 图片/PDF/Word/PPT/TXT"
        )

    contents = await file.read()
    if len(contents) > MAX_FILE_SIZE:
        raise HTTPException(status_code=400, detail="文件过大，最大10MB")

    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    safe_name = f"{timestamp}_{file.filename}"
    filepath = os.path.join(UPLOAD_DIR, safe_name)

    with open(filepath, "wb") as f:
        f.write(contents)

    extracted_text = extract_text_from_file(filepath, real_type, file.filename)

    return {
        "filename": file.filename,
        "content_type": real_type,
        "size": len(contents),
        "extracted_text": extracted_text[:5000],
        "file_path": f"/uploads/{safe_name}",
    }


@router.post("/files")
async def upload_multiple_files(
    files: list[UploadFile] = File(...),
    current_user: User = Depends(get_current_user),
):
    """批量上传文件并提取文本"""
    results = []
    all_text = []

    for file in files[:5]:
        ext = os.path.splitext(file.filename or "")[1].lower()
        real_type = _detect_file_type(file.filename or "", file.content_type)

        if real_type not in ALLOWED_TYPES and ext not in ALLOWED_EXTENSIONS:
            results.append({"filename": file.filename, "error": f"不支持的文件类型 ({ext})"})
            continue

        contents = await file.read()
        if len(contents) > MAX_FILE_SIZE:
            results.append({"filename": file.filename, "error": "文件过大"})
            continue

        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_name = f"{timestamp}_{file.filename}"
        filepath = os.path.join(UPLOAD_DIR, safe_name)

        with open(filepath, "wb") as f:
            f.write(contents)

        extracted = extract_text_from_file(filepath, real_type, file.filename)

        # 判断是否提取失败
        has_error = extracted.startswith("[") and ("失败" in extracted or "需要安装" in extracted or "旧版" in extracted)

        results.append({
            "filename": file.filename,
            "content_type": real_type,
            "size": len(contents),
            "extracted_text": extracted[:3000],
            "error": extracted if has_error else None,
        })

        if not has_error and extracted:
            all_text.append(f"【{file.filename}】\n{extracted}")

    return {
        "files": results,
        "combined_text": "\n\n".join(all_text)[:8000],
    }
